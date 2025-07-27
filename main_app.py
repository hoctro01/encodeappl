from ttkbootstrap import Window
from ttkbootstrap.constants import *
from ttkbootstrap import ttk
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext, simpledialog
import os
import json
import threading
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from encoder import get_resource_path

RULES_PATH = get_resource_path("rules.json")



class MainApp:
    def __init__(self):
        self.rules = []

        self.root = Window(themename="litera")  # 🟢 Phải tạo root trước
        self.root.title("Tool Mã Hoá / Giải Mã Văn Bản")
        self.root.geometry("900x700")

        # Gán icon .ico nếu có (dùng tuyệt đối để tránh lỗi path)
        icon_path = os.path.join(os.getcwd(), "app.ico")
        if os.path.exists(icon_path):
            self.root.iconbitmap(icon_path)

        self.var_mode = tk.StringVar(master=self.root, value="encode")  # ✅ Gán master
        self.theme_var = tk.StringVar(master=self.root, value="litera")  # ✅ Gán master

        self.nb = ttk.Notebook(self.root)
        self.nb.pack(fill="both", expand=True, padx=10, pady=10)

        self.build_main_tab()
        self.build_rules_tab()

        self.root.after(100, self.load_rules)
        self.root.mainloop()

    def build_main_tab(self):
        frame = ttk.Frame(self.nb)
        self.nb.add(frame, text="🔐 Encode / Decode")

        topbar = ttk.Frame(frame)
        topbar.pack(pady=10)
        ttk.Label(topbar, text="Chế độ:").pack(side="left", padx=5)
        ttk.Radiobutton(topbar, text="Mã hoá", variable=self.var_mode, value="encode").pack(side="left")
        ttk.Radiobutton(topbar, text="Giải mã", variable=self.var_mode, value="decode").pack(side="left")

        ttk.Button(frame, text="📂 Chọn File và Xử Lý", command=self.select_file).pack(pady=5)

        self.progress = ttk.Progressbar(frame, orient="horizontal", length=500, mode="determinate", maximum=100)
        self.progress.pack(pady=10)

        ttk.Label(frame, text="📄 Nội dung Gốc:").pack()
        self.text_before = scrolledtext.ScrolledText(frame, height=10)
        self.text_before.pack(fill="both", expand=True, padx=10)

        ttk.Label(frame, text="🛠️ Nội dung Sau Xử Lý:").pack()
        self.text_after = scrolledtext.ScrolledText(frame, height=10)
        self.text_after.pack(fill="both", expand=True, padx=10)

        ttk.Button(frame, text="🌓 Chuyển Theme", command=self.toggle_theme).pack(pady=5)

    def build_rules_tab(self):
        frame = ttk.Frame(self.nb)
        self.nb.add(frame, text="🛠️ Quản lý Rules")

        self.tree = ttk.Treeview(frame, columns=("from", "to", "enabled"), show="headings", height=15)
        self.tree.heading("from", text="Từ Gốc")
        self.tree.heading("to", text="Mã Hoá")
        self.tree.heading("enabled", text="Trạng Thái")
        self.tree.column("from", width=200)
        self.tree.column("to", width=200)
        self.tree.column("enabled", width=100)
        self.tree.pack(padx=10, pady=10, fill="both", expand=True)

        btns = ttk.Frame(frame)
        btns.pack(pady=5)
        ttk.Button(btns, text="➕ Thêm", command=self.add_rule).pack(side="left", padx=5)
        ttk.Button(btns, text="✏️ Sửa", command=self.edit_rule).pack(side="left", padx=5)
        ttk.Button(btns, text="❌ Xoá", command=self.delete_rule).pack(side="left", padx=5)
        ttk.Button(btns, text="🟢 Bật", command=self.enable_rule).pack(side="left", padx=5)
        ttk.Button(btns, text="🔴 Tắt", command=self.disable_rule).pack(side="left", padx=5)
        ttk.Button(btns, text="💾 Lưu Rules", command=self.save_rules).pack(side="left", padx=5)
        ttk.Button(btns, text="✅ Kiểm tra Rules", command=self.validate_rules).pack(side="left", padx=5)

    def toggle_theme(self):
        current = self.root.style.theme.name
        new_theme = "darkly" if current != "darkly" else "litera"
        self.root.style.theme_use(new_theme)

    def load_rules(self):
        if os.path.exists(RULES_PATH):
            with open(RULES_PATH, encoding='utf-8') as f:
                self.rules = json.load(f)
        self.reload_tree()

    def save_rules(self):
        try:
            with open(RULES_PATH, "w", encoding="utf-8") as f:
                json.dump(self.rules, f, ensure_ascii=False, indent=2)
            messagebox.showinfo("✅ Thành công", "Đã lưu rules.json!")
        except Exception as e:
            messagebox.showerror("❌ Lỗi khi lưu", str(e))

    def reload_tree(self):
        self.tree.delete(*self.tree.get_children())
        for rule in self.rules:
            self.tree.insert("", "end", values=(rule["from"], rule["to"], "✅" if rule.get("enabled", True) else "❌"))

    def get_selected_index(self):
        selection = self.tree.selection()
        if selection:
            return self.tree.index(selection[0])
        return None

    def add_rule(self):
        from_val = simpledialog.askstring("Từ gốc", "Nhập một từ gốc:", parent=self.root)
        if not from_val:
            return
        from_word = from_val.strip()
        all_from = set(r["from"] for r in self.rules)
        all_to = set(r["to"] for r in self.rules)
        if from_word in all_from:
            messagebox.showerror("Lỗi", f"Từ '{from_word}' đã tồn tại.")
            return
        index = 1
        while True:
            new_to = f"__ENC_{index:03}__"
            if new_to not in all_from and new_to not in all_to:
                break
            index += 1
        self.rules.append({"from": from_word, "to": new_to, "enabled": True})
        self.reload_tree()

    def edit_rule(self):
        idx = self.get_selected_index()
        if idx is None:
            return
        rule = self.rules[idx]
        new_from = simpledialog.askstring("Sửa từ gốc", "Nhập từ mới:", initialvalue=rule["from"])
        if not new_from:
            return
        new_word = new_from.strip()
        all_from = set(r["from"] for i, r in enumerate(self.rules) if i != idx)
        if new_word in all_from:
            messagebox.showerror("Lỗi", f"Từ '{new_word}' đã tồn tại.")
            return
        self.rules[idx]["from"] = new_word
        self.reload_tree()

    def delete_rule(self):
        idx = self.get_selected_index()
        if idx is not None:
            del self.rules[idx]
            self.reload_tree()

    def enable_rule(self):
        idx = self.get_selected_index()
        if idx is not None:
            self.rules[idx]["enabled"] = True
            self.reload_tree()

    def disable_rule(self):
        idx = self.get_selected_index()
        if idx is not None:
            self.rules[idx]["enabled"] = False
            self.reload_tree()

    def validate_rules(self):
        try:
            all_from = set()
            all_to = set()
            for rule in self.rules:
                f, t = rule["from"], rule["to"]
                if f in all_from:
                    raise ValueError(f"Trùng từ gốc: {f}")
                if t in all_to:
                    raise ValueError(f"Trùng mã hoá: {t}")
                if t in all_from or f == t:
                    raise ValueError(f"Mã hoá không hợp lệ: {t}")
                all_from.add(f)
                all_to.add(t)
            messagebox.showinfo("✅ OK", "Rules hợp lệ, không trùng lặp.")
        except Exception as e:
            messagebox.showerror("❌ Lỗi", str(e))

    def select_file(self):
        file_path = filedialog.askopenfilename(filetypes=[("Supported", "*.docx *.xlsx *.pptx *.txt")])
        if file_path:
            threading.Thread(target=self.process_file, args=(file_path,), daemon=True).start()

    def process_file(self, file_path):
        try:
            encode_map, decode_map = self.build_maps()
            rep_map = encode_map if self.var_mode.get() == "encode" else decode_map
            raw_text, new_text = "", ""
            ext = os.path.splitext(file_path)[1]

            self.progress["value"] = 0
            self.root.update_idletasks()

            if ext == ".docx":
                doc = Document(file_path)
                total_items = len(doc.paragraphs) + sum(len(t.rows) * len(t.columns) for t in doc.tables)
                count = 0
                for para in doc.paragraphs:
                    raw_text += para.text + "\n"
                    para.text = self.replace_text(para.text, rep_map)
                    count += 1
                    self.update_progress(count, total_items)
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            raw_text += cell.text + "\n"
                            cell.text = self.replace_text(cell.text, rep_map)
                            count += 1
                            self.update_progress(count, total_items)
                new_text = self.replace_text(raw_text, rep_map)
                doc.save(file_path.replace(".docx", f"_{self.var_mode.get()}.docx"))

            elif ext == ".xlsx":
                wb = load_workbook(file_path)
                total_cells = sum(len(row) for sheet in wb.worksheets for row in sheet.iter_rows())
                done_cells = 0
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows():
                        for cell in row:
                            if isinstance(cell.value, str):
                                raw_text += cell.value + "\n"
                                cell.value = self.replace_text(cell.value, rep_map)
                            done_cells += 1
                            self.update_progress(done_cells, total_cells)
                new_text = self.replace_text(raw_text, rep_map)
                wb.save(file_path.replace(".xlsx", f"_{self.var_mode.get()}.xlsx"))

            elif ext == ".pptx":
                prs = Presentation(file_path)
                total_shapes = sum(len(slide.shapes) for slide in prs.slides)
                count = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            raw_text += shape.text + "\n"
                            shape.text = self.replace_text(shape.text, rep_map)
                        count += 1
                        self.update_progress(count, total_shapes)
                new_text = self.replace_text(raw_text, rep_map)
                prs.save(file_path.replace(".pptx", f"_{self.var_mode.get()}.pptx"))

            elif ext == ".txt":
                with open(file_path, encoding="utf-8") as f:
                    raw_text = f.read()
                new_text = self.replace_text(raw_text, rep_map)
                with open(file_path.replace(".txt", f"_{self.var_mode.get()}.txt"), "w", encoding="utf-8") as f:
                    f.write(new_text)
                self.progress["value"] = 100

            self.text_before.delete(1.0, tk.END)
            self.text_before.insert(tk.END, raw_text.strip())

            self.text_after.delete(1.0, tk.END)
            self.text_after.insert(tk.END, new_text.strip())

            messagebox.showinfo("✅ Hoàn tất", f"Đã xử lý: {os.path.basename(file_path)}")
        except Exception as e:
            messagebox.showerror("❌ Lỗi", str(e))
        finally:
            self.progress["value"] = 0

    def update_progress(self, current, total):
        percent = int(current * 100 / total)
        self.progress["value"] = percent
        self.root.update_idletasks()

    def build_maps(self):
        encode_map = {}
        decode_map = {}
        for rule in self.rules:
            if not rule.get("enabled", True):
                continue
            encode_map[rule["from"]] = rule["to"]
            decode_map[rule["to"]] = rule["from"]
        return encode_map, decode_map

    def replace_text(self, text, rep_map):
        for k, v in rep_map.items():
            text = text.replace(k, v)
        return text


if __name__ == "__main__":
    MainApp()
